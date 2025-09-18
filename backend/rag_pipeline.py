# Cache and HF env fix - MUST be at very top BEFORE any transformers / sentence_transformers imports
import os
# Use writable cache directories to avoid PermissionError at /.cache in hosted environments
os.environ.setdefault("TRANSFORMERS_CACHE", "/tmp/transformers_cache")
os.environ.setdefault("HF_HOME", "/tmp/hf_home")
os.environ.setdefault("HUGGINGFACE_HUB_CACHE", "/tmp/hf_cache")
for d in (os.environ["TRANSFORMERS_CACHE"], os.environ["HF_HOME"], os.environ["HUGGINGFACE_HUB_CACHE"]):
    try:
        os.makedirs(d, exist_ok=True)
        os.chmod(d, 0o777)
    except Exception:
        pass

# ----- imports after cache env set -----
import uuid
from typing import List, Optional

import pandas as pd
import pytesseract
import torch
from docx import Document as DocxDocument
from dotenv import load_dotenv
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import Chroma
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation
from PyPDF2 import PdfReader
# Delay heavy transformer imports until needed
from transformers import AutoModelForSeq2SeqLM, AutoTokenizer, pipeline

load_dotenv()

# Config
EMBED_MODEL_NAME = os.getenv("EMBED_MODEL_NAME", "sentence-transformers/all-MiniLM-L6-v2")
GEN_MODEL = os.getenv("GEN_MODEL", "google/flan-t5-small")
CHROMA_DIR = os.getenv("CHROMA_DIR", "chroma_db")
COLLECTION_NAME = os.getenv("COLLECTION_NAME", "docs")

# Embeddings wrapper (will use TRANSFORMERS_CACHE set above)
hf_emb = HuggingFaceEmbeddings(model_name=EMBED_MODEL_NAME)

# Lazy generator pipeline
_gen_pipeline = None
def get_gen_pipeline():
    global _gen_pipeline
    if _gen_pipeline is None:
        cache_dir = os.environ.get("TRANSFORMERS_CACHE")
        # load tokenizer & model explicitly with cache_dir
        tokenizer = AutoTokenizer.from_pretrained(GEN_MODEL, cache_dir=cache_dir)
        model = AutoModelForSeq2SeqLM.from_pretrained(GEN_MODEL, cache_dir=cache_dir)
        device = 0 if torch.cuda.is_available() else -1
        _gen_pipeline = pipeline(
            "text2text-generation",
            model=model,
            tokenizer=tokenizer,
            device=device,
            truncation=True
        )
    return _gen_pipeline

# Prompt templates
PROMPT_TEMPLATES = {
    "pharma": "You are a domain expert in pharmaceutical documentation. Answer concisely and cite sources when possible.",
    "legal": "You are a legal assistant. Provide careful, precise answers and state limits if unsure.",
    "engineering": "You are an engineering assistant. Answer with clarity and refer to the provided documents.",
    "general": "You are a helpful assistant. Use the provided context to answer the question and cite sources."
}

def build_system_prompt(domain: Optional[str]) -> str:
    return PROMPT_TEMPLATES.get((domain or "").lower(), PROMPT_TEMPLATES["general"])

# Extractors (same as before)
def extract_text_from_pdf(path: str) -> str:
    reader = PdfReader(path)
    texts = []
    for page in reader.pages:
        t = page.extract_text()
        if t:
            texts.append(t)
    text = "\n".join(texts)
    # fallback to OCR if no text found
    if not text.strip():
        try:
            pages = convert_from_path(path)
            ocr_pages = [pytesseract.image_to_string(page) for page in pages]
            text = "\n".join(ocr_pages)
        except Exception:
            pass
    return text

def extract_text_from_docx(path: str) -> str:
    doc = DocxDocument(path)
    return "\n".join([p.text for p in doc.paragraphs if p.text])

def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shp in slide.shapes:
            if hasattr(shp, "text"):
                texts.append(shp.text)
    return "\n".join(texts)

def extract_text_from_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def extract_text_from_csv(path: str, max_rows=200) -> str:
    try:
        df = pd.read_csv(path, nrows=max_rows)
        return df.to_string()
    except Exception:
        return extract_text_from_txt(path)

def ocr_image(path: str) -> str:
    img = Image.open(path).convert("RGB")
    return pytesseract.image_to_string(img)

def extract_text_from_file(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(path)
    elif ext == ".docx":
        return extract_text_from_docx(path)
    elif ext in [".pptx"]:
        return extract_text_from_pptx(path)
    elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
        return ocr_image(path)
    elif ext == ".csv":
        return extract_text_from_csv(path)
    elif ext == ".txt":
        return extract_text_from_txt(path)
    else:
        return extract_text_from_txt(path)

# Chunking + ingestion
def chunk_text(text: str, chunk_size:int=800, chunk_overlap:int=150) -> List[str]:
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return splitter.split_text(text)

def ingest_document(path: str, doc_id: Optional[str]=None, domain: Optional[str]=None, tags: Optional[str]=None) -> int:
    os.makedirs(CHROMA_DIR, exist_ok=True)
    text = extract_text_from_file(path)
    chunks = chunk_text(text)
    metadatas = []
    for i in range(len(chunks)):
        md = {"source": os.path.basename(path), "chunk_index": i}
        if domain: md["domain"] = domain
        if tags: md["tags"] = tags
        metadatas.append(md)
    vectordb = Chroma.from_texts(
        texts=chunks,
        embedding=hf_emb,
        metadatas=metadatas,
        persist_directory=CHROMA_DIR,
        collection_name=COLLECTION_NAME
    )
    vectordb.persist()
    return len(chunks)

def build_retriever(k:int=5, domain: Optional[str]=None):
    vectordb = Chroma(persist_directory=CHROMA_DIR, embedding_function=hf_emb, collection_name=COLLECTION_NAME)
    retriever = vectordb.as_retriever(search_kwargs={"k": k})
    # simple domain filter (Chroma supports metadata filtering via where in some versions)
    return retriever

# Generation
def generate_answer_from_context(context: str, question: str, domain: Optional[str]=None, max_length:int=256) -> str:
    gen = get_gen_pipeline()
    system = build_system_prompt(domain)
    prompt = f"{system}\n\nContext:\n{context}\n\nQuestion: {question}\nAnswer:"
    out = gen(prompt, max_length=max_length, do_sample=False)[0]["generated_text"]
    return out

def answer_query_local(question: str, retriever, top_k: int=5, domain: Optional[str]=None):
    docs = retriever.get_relevant_documents(question)
    contexts = [d.page_content for d in docs][:top_k]
    combined = "\n\n".join(contexts)
    answer = generate_answer_from_context(combined, question, domain=domain)
    citations = [{"source": d.metadata.get("source","unknown"), "chunk_index": d.metadata.get("chunk_index",-1), "text": d.page_content[:500]} for d in docs[:top_k]]
    return {"answer": answer, "citations": citations}
